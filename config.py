import os
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_experimental.text_splitter import SemanticChunker
from sentence_transformers import CrossEncoder
from langchain_community.document_loaders import PyPDFLoader, CSVLoader
from langchain_core.documents import Document
from google import genai
from dotenv import load_dotenv

load_dotenv()

client = genai.Client()
cross_encoder = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2", device="cpu")

COLLECTION_NAME = "docutrust"
EMBED_MODEL = "BAAI/bge-small-en-v1.5"

embeddings = HuggingFaceEmbeddings(model_name=EMBED_MODEL, model_kwargs={"device": "cpu"})
text_splitter = SemanticChunker(embeddings, breakpoint_threshold_type="percentile")


def _make_fresh_store() -> Chroma:
    import chromadb
    raw = chromadb.EphemeralClient()
    return Chroma(client=raw, collection_name=COLLECTION_NAME, embedding_function=embeddings)


vector_store = _make_fresh_store()


def _table_to_text(table) -> str:
    """
    Convert a python-docx or pptx table into a readable key:value text block.
    Each row becomes 'Col1: val1 | Col2: val2 | ...' so the LLM can reason over it.
    If the first row looks like a header, it is used as column labels for all rows.
    """
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    # Detect header row — use it to label subsequent rows
    header = rows[0]
    all_empty_header = all(h == "" for h in header)

    lines = []
    if not all_empty_header and len(rows) > 1:
        # Use header labels
        for data_row in rows[1:]:
            if any(c for c in data_row):
                parts = [f"{h}: {v}" for h, v in zip(header, data_row) if h or v]
                lines.append(" | ".join(parts))
        # Also include the raw header row for context
        lines.insert(0, "Columns: " + " | ".join(header))
    else:
        # No clear header — just join cells per row
        for row in rows:
            if any(c for c in row):
                lines.append(" | ".join(c for c in row if c))

    return "\n".join(lines)


def load_file(file_path: str):
    """
    Load any supported file with full table awareness.
    Tables are converted to structured key:value text so the LLM
    can reason over tabular data correctly.
    """
    ext = os.path.splitext(file_path)[1].lower()

    # ── PDF ──────────────────────────────────────────────────────────────
    # PDFs with tables: pdfplumber extracts tables as structured text.
    # Falls back to PyPDFLoader if pdfplumber is not installed.
    if ext == ".pdf":
        try:
            import pdfplumber
            docs = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    parts = []
                    # Extract tables first, structured as key:value rows
                    for table in page.extract_tables():
                        if not table:
                            continue
                        header = [str(c).strip() if c else "" for c in table[0]]
                        all_empty = all(h == "" for h in header)
                        if not all_empty and len(table) > 1:
                            parts.append("Columns: " + " | ".join(header))
                            for row in table[1:]:
                                cells = [str(c).strip() if c else "" for c in row]
                                if any(cells):
                                    parts.append(" | ".join(
                                        f"{h}: {v}" for h, v in zip(header, cells) if h or v
                                    ))
                        else:
                            for row in table:
                                cells = [str(c).strip() if c else "" for c in row]
                                if any(cells):
                                    parts.append(" | ".join(c for c in cells if c))
                    # Extract remaining text (non-table)
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        parts.append(page_text)
                    if parts:
                        docs.append(Document(
                            page_content="\n".join(parts),
                            metadata={"source": file_path, "page": i}
                        ))
            if docs:
                return docs
            # pdfplumber found nothing — fall through to PyPDFLoader
        except ImportError:
            pass  # pdfplumber not installed — use PyPDFLoader below
        except Exception:
            pass  # corrupted or encrypted PDF — fall through
        try:
            return PyPDFLoader(file_path).load()
        except Exception as e:
            raise RuntimeError(f"PDF failed: {e}") from e

    # ── CSV ──────────────────────────────────────────────────────────────
    # CSV is already tabular — CSVLoader converts each row to a Document
    # with 'column: value' format natively.
    elif ext == ".csv":
        try:
            return CSVLoader(file_path, encoding="utf-8").load()
        except Exception as e:
            raise RuntimeError(f"CSV failed: {e}") from e

    # ── DOCX / DOC ───────────────────────────────────────────────────────
    # Extract paragraphs AND tables in document order.
    elif ext in (".docx", ".doc"):
        try:
            from docx import Document as DocxDoc
            from docx.oxml.ns import qn
            doc = DocxDoc(file_path)
            parts = []

            # Iterate body elements in order to preserve paragraph/table sequence
            for block in doc.element.body:
                tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

                if tag == "p":
                    # Paragraph
                    text = block.text_content() if hasattr(block, "text_content") else ""
                    # Fallback: join all 't' (run text) elements
                    if not text:
                        text = "".join(t.text or "" for t in block.iter()
                                       if t.tag.endswith("}t"))
                    if text.strip():
                        parts.append(text.strip())

                elif tag == "tbl":
                    # Table — find the matching python-docx Table object by index
                    # We re-iterate tables to match by XML element position
                    pass  # handled below

            # Simpler approach: interleave paragraphs and tables by re-parsing
            parts = []
            body_elements = list(doc.element.body)
            tbl_idx = 0
            para_idx = 0
            all_tables = doc.tables
            all_paras  = doc.paragraphs

            for block in body_elements:
                tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
                if tag == "p":
                    if para_idx < len(all_paras):
                        txt = all_paras[para_idx].text.strip()
                        if txt:
                            parts.append(txt)
                        para_idx += 1
                elif tag == "tbl":
                    if tbl_idx < len(all_tables):
                        tbl_text = _table_to_text(all_tables[tbl_idx])
                        if tbl_text:
                            parts.append("[TABLE]\n" + tbl_text + "\n[/TABLE]")
                        tbl_idx += 1

            if not parts:
                raise RuntimeError("No content found in DOCX.")
            return [Document(
                page_content="\n".join(parts),
                metadata={"source": file_path, "page": 1}
            )]
        except ImportError:
            raise RuntimeError("DOCX failed: run pip install python-docx")
        except RuntimeError:
            raise
        except Exception as e:
            raise RuntimeError(f"DOCX failed: {e}") from e

    # ── PPTX / PPT ───────────────────────────────────────────────────────
    # Extract text AND tables from each slide.
    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            from pptx.util import Pt
            prs = Presentation(file_path)
            docs = []
            for i, slide in enumerate(prs.slides, start=1):
                parts = []
                for shape in slide.shapes:
                    # Text frames
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            parts.append(txt)
                    # Tables inside slides
                    if shape.has_table:
                        tbl_text = _table_to_text(shape.table)
                        if tbl_text:
                            parts.append("[TABLE]\n" + tbl_text + "\n[/TABLE]")
                if parts:
                    docs.append(Document(
                        page_content="\n".join(parts),
                        metadata={"source": file_path, "page": i}
                    ))
            return docs
        except ImportError:
            raise RuntimeError("PPTX failed: run pip install python-pptx")
        except Exception as e:
            raise RuntimeError(f"PPTX failed: {e}") from e

    # ── MARKDOWN ─────────────────────────────────────────────────────────
    # Markdown tables are already plain text with | separators — kept as-is.
    elif ext == ".md":
        try:
            with open(file_path, encoding="utf-8", errors="replace") as fh:
                text = fh.read()
            return [Document(page_content=text, metadata={"source": file_path, "page": 1})]
        except Exception as e:
            raise RuntimeError(f"Markdown failed: {e}") from e

    # ── HTML / HTM ───────────────────────────────────────────────────────
    # Parse HTML tables into structured key:value text using BeautifulSoup.
    elif ext in (".html", ".htm"):
        try:
            from bs4 import BeautifulSoup
            with open(file_path, encoding="utf-8", errors="replace") as fh:
                soup = BeautifulSoup(fh, "html.parser")

            parts = []
            for element in soup.find_all(["p", "h1", "h2", "h3", "h4",
                                           "h5", "h6", "li", "table"]):
                if element.name == "table":
                    # Extract HTML table rows
                    rows = element.find_all("tr")
                    if not rows:
                        continue
                    header_cells = rows[0].find_all(["th", "td"])
                    header = [c.get_text(strip=True) for c in header_cells]
                    all_empty = all(h == "" for h in header)

                    if not all_empty and len(rows) > 1:
                        parts.append("[TABLE]")
                        parts.append("Columns: " + " | ".join(header))
                        for row in rows[1:]:
                            cells = [c.get_text(strip=True)
                                     for c in row.find_all(["th", "td"])]
                            if any(cells):
                                parts.append(" | ".join(
                                    f"{h}: {v}" for h, v in zip(header, cells) if h or v
                                ))
                        parts.append("[/TABLE]")
                    else:
                        for row in rows:
                            cells = [c.get_text(strip=True)
                                     for c in row.find_all(["th", "td"])]
                            if any(cells):
                                parts.append(" | ".join(c for c in cells if c))
                else:
                    txt = element.get_text(strip=True)
                    if txt:
                        parts.append(txt)

            text = "\n".join(parts)
            if not text:
                raise RuntimeError("No text extracted from HTML.")
            return [Document(page_content=text, metadata={"source": file_path, "page": 1})]
        except ImportError:
            raise RuntimeError("HTML failed: run pip install beautifulsoup4")
        except RuntimeError:
            raise
        except Exception as e:
            raise RuntimeError(f"HTML failed: {e}") from e

    # ── UNSUPPORTED ──────────────────────────────────────────────────────
    else:
        raise RuntimeError(
            f"Unsupported file type: '{ext}'. Supported: pdf, csv, docx, pptx, md, html"
        )


def _clear_documents_folder():
    doc_dir = "./documents"
    if os.path.exists(doc_dir):
        for fname in os.listdir(doc_dir):
            try:
                os.remove(os.path.join(doc_dir, fname))
            except Exception as e:
                print(f"  Could not delete {fname}: {e}")
    os.makedirs(doc_dir, exist_ok=True)
    print("  Cleared ./documents folder on startup.")


_clear_documents_folder()